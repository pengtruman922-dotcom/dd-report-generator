"""Extract text from PowerPoint (.pptx) files.

Extracts text from text frames, tables, and embedded images (via OCR).
Handles grouped shapes recursively.
"""

from __future__ import annotations

import logging
from pathlib import Path

from parsers.ocr_utils import ocr_image

log = logging.getLogger(__name__)


def extract_pptx_text(file_path) -> str:
    """Return the full text of a .pptx file, including OCR of embedded images.

    Args:
        file_path: A file path (str/Path) or a file-like object (e.g. BytesIO).
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        log.warning("python-pptx not installed, cannot parse .pptx")
        return ""

    try:
        # python-pptx Presentation() accepts both file paths and file-like objects
        if isinstance(file_path, (str, Path)):
            prs = Presentation(str(file_path))
        else:
            prs = Presentation(file_path)
    except Exception:
        log.exception("Failed to open .pptx: %s", file_path)
        return ""

    slide_parts: list[str] = []
    total_images = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        ocr_texts: list[str] = []

        def _process_shape(shape):
            nonlocal total_images

            # Text frames (text boxes, titles, etc.)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    row_str = " | ".join(cells)
                    if row_str.replace("|", "").strip():
                        texts.append(row_str)

            # Images → OCR
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    ocr_text = ocr_image(img_bytes)
                    if ocr_text.strip():
                        ocr_texts.append(ocr_text)
                        total_images += 1
                except Exception:
                    pass

            # Recurse into group shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    _process_shape(child)

        for shape in slide.shapes:
            _process_shape(shape)

        # Combine text + OCR for this slide
        slide_content = texts + ocr_texts
        if slide_content:
            slide_parts.append(f"[Slide {slide_idx}]\n" + "\n".join(slide_content))

    if total_images > 0:
        log.info("PPTX: extracted text from %d slides + %d images (OCR) in %s",
                 len(prs.slides), total_images, file_path)

    return "\n\n".join(slide_parts)
